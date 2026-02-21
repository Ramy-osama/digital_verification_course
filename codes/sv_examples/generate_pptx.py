"""
Generate Section2_arrays.pptx - LRM Sections 3-7 with QuestaSim Introduction
Matches Ain Shams University template style:
  - Blue header band with university branding
  - Footer: date, "Ain Shams University - Faculty of Engineering", slide number
  - Code slides with monospaced font, syntax highlighting, and expected output

Usage:
    pip install python-pptx
    python generate_pptx.py
"""

import os
import re
import datetime
import textwrap
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BLUE = RGBColor(0x1F, 0x4E, 0x79)
DARK_BLUE = RGBColor(0x0D, 0x2E, 0x4E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
BLACK = RGBColor(0x00, 0x00, 0x00)
CODE_BG = RGBColor(0xF5, 0xF5, 0xF5)
OUTPUT_GREEN = RGBColor(0x00, 0x64, 0x00)

KW_BLUE = RGBColor(0x00, 0x00, 0xCC)
SYS_TEAL = RGBColor(0x00, 0x80, 0x80)
COMMENT_GREEN = RGBColor(0x00, 0x80, 0x00)
LITERAL_PURPLE = RGBColor(0x80, 0x00, 0x80)

ACCENT_ORANGE = RGBColor(0xC0, 0x50, 0x20)
MEDIUM_GRAY = RGBColor(0x55, 0x55, 0x55)
SECTION_TEAL = RGBColor(0x00, 0x6B, 0x6B)
STRING_BROWN = RGBColor(0xA3, 0x11, 0x15)

SV_KEYWORDS = {
    "module", "endmodule", "struct", "typedef", "packed", "signed",
    "unsigned", "union", "tagged", "int", "bit", "logic", "byte",
    "shortint", "longint", "integer", "time", "real", "shortreal",
    "string", "void", "reg", "event", "chandle", "enum", "const",
    "static", "automatic", "rand", "randc", "new", "null", "this",
    "super", "class", "endclass", "function", "endfunction", "task",
    "endtask", "begin", "end", "if", "else", "for", "foreach",
    "while", "do", "return", "initial", "always", "assign", "ref",
    "input", "output", "inout", "extends", "virtual", "pure",
    "import", "export", "constraint", "with", "inside", "default",
    "wire", "supply0", "supply1", "tri", "triand", "trior", "tri0",
    "tri1", "wand", "wor", "always_ff", "always_comb", "posedge",
    "negedge", "interface", "endinterface", "modport", "package",
    "endpackage", "program", "endprogram", "parameter", "localparam",
    "case", "endcase", "timeunit", "timeprecision",
}

_TOKEN_RE = re.compile(
    r"(//.*)"
    r"|(\$\w+)"
    r"|(\d+'[bBhHdDoO][0-9a-fA-F_xXzZ]+)"
    r'|("[^"]*")'
    r"|(\b[a-zA-Z_]\w*\b)"
    r"|(\S)"
    r"|( +)"
)


def tokenize_sv_line(line):
    """Split a SystemVerilog line into (text, category) tokens."""
    tokens = []
    pos = 0
    for m in _TOKEN_RE.finditer(line):
        if m.start() > pos:
            tokens.append((line[pos:m.start()], "default"))
        comment, sysfn, literal, strlit, ident, punct, space = m.groups()
        if comment is not None:
            tokens.append((comment, "comment"))
        elif sysfn is not None:
            tokens.append((sysfn, "system"))
        elif literal is not None:
            tokens.append((literal, "literal"))
        elif strlit is not None:
            tokens.append((strlit, "string"))
        elif ident is not None:
            if ident in SV_KEYWORDS:
                tokens.append((ident, "keyword"))
            else:
                tokens.append((ident, "default"))
        elif punct is not None:
            tokens.append((punct, "default"))
        elif space is not None:
            tokens.append((space, "default"))
        pos = m.end()
    if pos < len(line):
        tokens.append((line[pos:], "default"))
    return tokens


TCL_KEYWORDS = {
    "vlib", "vlog", "vcom", "vsim", "vopt", "vmap", "vcover",
    "run", "quit", "add", "wave", "force", "examine", "do",
    "transcript", "log", "bp", "describe", "drivers", "step",
    "restart", "checkpoint", "restore", "show", "radix",
    "onbreak", "onerror", "when", "quietly", "echo", "cd",
}

_TCL_TOKEN_RE = re.compile(
    r"(#.*)"
    r'|("[^"]*")'
    r"|(-\w+)"
    r"|(\b[a-zA-Z_]\w*\b)"
    r"|(\S)"
    r"|( +)"
)


def tokenize_tcl_line(line):
    """Split a Tcl/do-file line into (text, category) tokens."""
    tokens = []
    pos = 0
    for m in _TCL_TOKEN_RE.finditer(line):
        if m.start() > pos:
            tokens.append((line[pos:m.start()], "default"))
        comment, strlit, flag, ident, punct, space = m.groups()
        if comment is not None:
            tokens.append((comment, "comment"))
        elif strlit is not None:
            tokens.append((strlit, "string"))
        elif flag is not None:
            tokens.append((flag, "literal"))
        elif ident is not None:
            if ident in TCL_KEYWORDS:
                tokens.append((ident, "keyword"))
            else:
                tokens.append((ident, "default"))
        elif punct is not None:
            tokens.append((punct, "default"))
        elif space is not None:
            tokens.append((space, "default"))
        pos = m.end()
    if pos < len(line):
        tokens.append((line[pos:], "default"))
    return tokens


TOKEN_STYLES = {
    "keyword": (KW_BLUE, True, False),
    "system":  (SYS_TEAL, False, False),
    "comment": (COMMENT_GREEN, False, True),
    "literal": (LITERAL_PURPLE, False, False),
    "string":  (STRING_BROWN, False, False),
    "default": (BLACK, False, False),
}

SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(7.5)

TODAY = datetime.date.today().strftime("%m/%d/%Y")

OUTPUT_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                          "..", "..", "Presentations")
OUTPUT_FILE = os.path.join(OUTPUT_DIR, "Section2.pptx")


def add_header_band(slide, title_text):
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        SLIDE_WIDTH, Inches(1.0)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = BLUE
    header.line.fill.background()

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.0),
        SLIDE_WIDTH, Inches(0.06)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(0xC0, 0x50, 0x20)
    accent.line.fill.background()

    txBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.15), Inches(9), Inches(0.7)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"


def add_footer(slide, slide_num):
    left_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(7.0), Inches(1.5), Inches(0.4)
    )
    tf = left_box.text_frame
    p = tf.paragraphs[0]
    p.text = TODAY
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
    p.font.name = "Calibri"

    center_box = slide.shapes.add_textbox(
        Inches(2.5), Inches(7.0), Inches(5), Inches(0.4)
    )
    tf = center_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Ain Shams University - Faculty of Engineering"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
    p.font.name = "Calibri"
    p.font.italic = True

    right_box = slide.shapes.add_textbox(
        Inches(8.5), Inches(7.0), Inches(1.2), Inches(0.4)
    )
    tf = right_box.text_frame
    p = tf.paragraphs[0]
    p.text = str(slide_num)
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
    p.font.name = "Calibri"


def _add_styled_run(paragraph, text, category, font_size=Pt(11)):
    """Add a single run with syntax-highlighting style."""
    color, bold, italic = TOKEN_STYLES.get(category, (BLACK, False, False))
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Consolas"
    run.font.size = font_size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def add_code_box(slide, code_text, top, height, output_text=None, lang="sv"):
    code_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.4), top, Inches(9.2), height
    )
    code_shape.fill.solid()
    code_shape.fill.fore_color.rgb = CODE_BG
    code_shape.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    code_shape.line.width = Pt(1)

    tf = code_shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)

    tokenizer = tokenize_tcl_line if lang == "tcl" else tokenize_sv_line
    cleaned = textwrap.dedent(code_text).strip()

    for i, line in enumerate(cleaned.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        p.space_after = Pt(1)
        p.space_before = Pt(0)
        pPr = p._p.get_or_add_pPr()
        pPr.set('indent', '0')
        pPr.set('marL', '0')

        tokens = tokenizer(line)
        if not tokens:
            _add_styled_run(p, "", "default")
        else:
            for text, category in tokens:
                _add_styled_run(p, text, category)

    if output_text:
        p = tf.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = "// Simulation Output:"
        run.font.name = "Consolas"
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = OUTPUT_GREEN

        for line in output_text.strip().split("\n"):
            p = tf.add_paragraph()
            p.space_after = Pt(1)
            p.space_before = Pt(0)
            run = p.add_run()
            run.text = "// " + line
            run.font.name = "Consolas"
            run.font.size = Pt(10)
            run.font.color.rgb = OUTPUT_GREEN


def _classify_bullet(text):
    """Return (display_text, color, bold, size) based on line content."""
    stripped = text.lstrip()
    indent = len(text) - len(stripped)

    if stripped == "":
        return ("", DARK_BLUE, False, Pt(8))

    is_numbered = bool(re.match(r'^\d+\.', stripped))
    is_section_num = bool(re.match(r'^\d+\.\d+', stripped))
    is_sub = indent >= 4

    if is_numbered:
        return (stripped, ACCENT_ORANGE, True, Pt(17))
    elif is_section_num:
        return (stripped, BLUE, True, Pt(17))
    elif is_sub:
        return (stripped, MEDIUM_GRAY, False, Pt(15))
    else:
        return (stripped, DARK_BLUE, False, Pt(17))


def add_bullet_slide(prs, title, bullets, slide_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_band(slide, title)
    add_footer(slide, slide_num)

    txBox = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.3), Inches(8.8), Inches(5.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        display, color, bold, size = _classify_bullet(bullet)

        if display == "":
            p.space_after = Pt(2)
            p.space_before = Pt(0)
            run = p.add_run()
            run.text = ""
            run.font.size = Pt(4)
        else:
            is_sub = bullet.lstrip() != bullet and (len(bullet) - len(bullet.lstrip())) >= 4
            p.space_after = Pt(4)
            p.space_before = Pt(0)
            if is_sub:
                pPr = p._p.get_or_add_pPr()
                pPr.set('marL', str(Emu(Inches(0.4))))
            run = p.add_run()
            run.text = display
            run.font.name = "Calibri"
            run.font.size = size
            run.font.color.rgb = color
            run.font.bold = bold

    return slide


def add_code_slide(prs, title, code, output, slide_num, lang="sv"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_band(slide, title)
    add_footer(slide, slide_num)
    add_code_box(slide, code, Inches(1.3), Inches(5.5), output, lang=lang)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    sn = 0

    # ================================================================
    #  TITLE SLIDE
    # ================================================================
    sn += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = BLUE
    bg.line.fill.background()

    txBox = slide.shapes.add_textbox(
        Inches(1), Inches(2.2), Inches(8), Inches(2.0)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Section 2"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"

    p2 = tf.add_paragraph()
    p2.space_before = Pt(40)
    p2.text = "Ramy Osama"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(28)
    p2.font.color.rgb = RGBColor(0xCC, 0xDD, 0xFF)
    p2.font.name = "Calibri"

    footer_box = slide.shapes.add_textbox(
        Inches(1), Inches(5.5), Inches(8), Inches(0.8)
    )
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Ain Shams University - Faculty of Engineering"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0xAA, 0xCC, 0xFF)
    p.font.name = "Calibri"
    p.font.italic = True

    # ================================================================
    #  COMPILATION, ELABORATION & SIMULATION PHASES
    # ================================================================
    sn += 1
    add_bullet_slide(prs, "Compilation, Elaboration & Simulation", [
        "1. Compilation  (vlog -sv design.sv)",
        "    - Parses each source file independently",
        "    - Checks syntax and language rules (IEEE 1800)",
        "    - Converts source to an intermediate representation",
        "    - Errors: missing semicolons, undeclared variables, bad syntax",
        "",
        "2. Elaboration  (vsim work.top_module)",
        "    - Resolves the full module hierarchy and instantiation tree",
        "    - Evaluates parameters, generates parameterized instances",
        "    - Connects ports, checks width mismatches and unconnected ports",
        "    - Errors: missing modules, port count/type mismatches",
        "",
        "3. Simulation  (run -all)",
        "    - Executes the event-driven simulation model",
        "    - Processes initial, always, always_ff, always_comb blocks",
        "    - Advances simulation time, schedules events",
        "    - Errors: assertion failures, timeout, $fatal, x-propagation",
    ], sn)

    sn += 1
    add_bullet_slide(prs, "Simulator Flow - Phase by Phase", [
        "Source (.sv)  -->  vlog  -->  vsim  -->  run  -->  Results",
        "",
        "Phase 1 - COMPILATION (vlog):",
        "    .sv files --> syntax check --> intermediate objects in 'work' library",
        "",
        "Phase 2 - ELABORATION (vsim):",
        "    Top module --> resolve hierarchy --> expand parameters --> bind ports",
        "",
        "Phase 3 - SIMULATION (run):",
        "    Time 0 --> execute events --> advance time --> repeat until $finish",
        "",
        "Key insight: A file can compile successfully but fail at elaboration",
        "  (e.g. instantiating a module with wrong port names),",
        "  or pass elaboration but fail at simulation (e.g. runtime assertion).",
    ], sn)

    # ================================================================
    #  QUESTASIM INTRODUCTION
    # ================================================================
    sn += 1
    add_bullet_slide(prs, "What is QuestaSim?", [
        "Industry-standard HDL simulator by Siemens EDA (formerly Mentor Graphics)",
        "Supports Verilog, SystemVerilog, VHDL, and SystemC",
        "Used for RTL simulation, functional verification, and design debug",
        "Provides GUI mode, command-line mode, and batch mode",
        "IEEE compliant: 1800 (SystemVerilog), 1364 (Verilog), 1076 (VHDL)",
        "Includes code coverage, assertions, and waveform analysis",
    ], sn)

    sn += 1
    add_bullet_slide(prs, "QuestaSim - Basic Simulation Flow", [
        "Step 1:  Create a working library          (vlib work)",
        "Step 2:  Map the library                          (vmap work work)",
        "Step 3:  Compile the design                   (vlog -sv design.sv)",
        "Step 4:  Load the design into simulator  (vsim work.top_module)",
        "Step 5:  Run the simulation                    (run -all)",
        "Step 6:  Debug with waveforms              (add wave, force)",
    ], sn)

    sn += 1
    add_code_slide(prs, "Step 1-2: Create Library & Compile", """
# Step 1: Create a working library
vlib work

# Step 2: Map logical name to physical directory
vmap work work

# Step 3: Compile Verilog files
vlog my_design.v

# Compile SystemVerilog files (use -sv flag)
vlog -sv my_testbench.sv

# Compile multiple files at once
vlog -sv file1.sv file2.sv file3.sv
""", None, sn, lang="tcl")

    sn += 1
    add_code_slide(prs, "Step 3-4: Load Design & Run Simulation", """
# Load the top-level module for simulation
vsim work.testbench

# Enable full signal access for debugging
vsim -voptargs=+acc work.testbench

# Add all signals to the wave window
add wave -r /*

# Run simulation for a specific time
run 1000ns

# Run until $finish is encountered
run -all
""", None, sn, lang="tcl")

    sn += 1
    add_code_slide(prs, "Step 5: Debugging Commands", """
# Force a signal to a specific value
force /testbench/rst 1 0, 0 100ns

# Log signals for waveform recording
log -r /*

# Restart simulation from time 0
restart -f

# Step through simulation
step
""", None, sn, lang="tcl")

    sn += 1
    add_code_slide(prs, "Complete Example: sim.do Script", """
# sim.do - Complete simulation script
# Run with: do sim.do (in QuestaSim transcript)

vlib work
vlog -sv design.sv testbench.sv

vsim -voptargs=+acc work.testbench

add wave -r /*
run -all

quit -sim
""", None, sn, lang="tcl")

    sn += 1
    add_bullet_slide(prs, "QuestaSim - Modes of Operation", [
        "GUI Mode (default): Interactive windows, menus, and waveform viewer",
        "    Invoke with:  vsim",
        "Command-line Mode: Interactive transcript, no graphical windows",
        "    Invoke with:  vsim -c",
        "Batch Mode: Non-interactive, scripts only, for regression runs",
        "    Invoke with:  vsim -c -do \"do sim.do; quit -f\"",
        "Use .do files (macro scripts) to automate repetitive tasks",
    ], sn)

    sn += 1
    add_bullet_slide(prs, "QuestaSim - Key Tips for Beginners", [
        "Always create the library (vlib work) before compiling",
        "Use -sv flag with vlog for SystemVerilog features",
        "Use -voptargs=+acc with vsim for full signal visibility",
        "Save commands in .do files for repeatable simulations",
        "Use 'add wave -r /*' to quickly view all signals",
        "Check the transcript window for compilation errors",
        "Use $display and $monitor in testbenches for text-based debug",
    ], sn)

    # ================================================================
    #  SECTION 5: LEXICAL CONVENTIONS
    # ================================================================
    sn += 1
    add_bullet_slide(prs, "Section 5 - Lexical Conventions", [
        "5.7   Number formats: sized/unsized, bases (bin, oct, dec, hex)",
        "5.7   Special values: 'x (unknown), 'z (high-impedance), '0, '1",
        "5.7   Signed literals: 4'shF = -1 (signed hex)",
        "5.8   String literals and string methods",
        "5.11  Operators: arithmetic, logical, bitwise, reduction, shift",
        "5.11  Concatenation {a,b}, replication {4{x}}, conditional (? :)",
    ], sn)

    sn += 1
    add_code_slide(prs, "5.7 Number Formats & Literals", """
// Sized literals
$display("4'b1001   = %b (decimal %0d)", 4'b1001, 4'b1001);
$display("8'hFF     = %h (decimal %0d)", 8'hFF, 8'hFF);
$display("12'o7460  = %o (decimal %0d)", 12'o7460, 12'o7460);

// Signed literals
$display("4'shF     = %0d (signed)", $signed(4'shF));

// Special values
logic [7:0] v;
v = 8'hxF;   $display("8'hxF  = %h", v);
v = '1;      $display("'1     = %b (all ones)", v);
v = '0;      $display("'0     = %b (all zeros)", v);

// Real and string
real r = 3.14;
string s = "Hello, SV!";
$display("real = %f, string = %s (len=%0d)", r, s, s.len());
""",
    "4'b1001 = 1001 (decimal 9)\n"
    "8'hFF = ff (decimal 255)\n"
    "4'shF = -1 (signed)\n"
    "'1 = 11111111, '0 = 00000000", sn)

    sn += 1
    add_code_slide(prs, "5.11 Operators Overview", """
logic [7:0] a = 8'hA5, b = 8'h3C;

// Arithmetic
$display("%h + %h = %h", a, b, a + b);
$display("%h * 2  = %h", a, a * 2);

// Bitwise
$display("%h & %h = %h", a, b, a & b);
$display("%h | %h = %h", a, b, a | b);
$display("~%h     = %h", a, ~a);

// Reduction
$display("&FF = %b, |00 = %b", &8'hFF, |8'h00);

// Shift
$display("A5 << 2 = %b", a << 2);

// Concatenation & replication
$display("{AB,CD} = %h", {8'hAB, 8'hCD});
$display("{4{2'b10}} = %b", {4{2'b10}});

// Conditional
$display("1 ? AA : 55 = %h", 1 ? 8'hAA : 8'h55);
""",
    "a5 + 3c = e1, a5 * 2 = 4a\n"
    "a5 & 3c = 24, a5 | 3c = bd, ~a5 = 5a\n"
    "&FF = 1, |00 = 0\n"
    "{AB,CD} = abcd, {4{10}} = 10101010\n"
    "1 ? AA : 55 = aa", sn)

    # ================================================================
    #  LAB: SECTION 5
    # ================================================================
    sn += 1
    add_bullet_slide(prs, "Lab Exercise - Lexical Conventions", [
        "Display the value 255 in binary, octal, decimal, and hex",
        "using $display with format specifiers %b, %o, %d, %h.",
    ], sn)

    # ================================================================
    #  SECTION 6: DATA TYPES
    # ================================================================
    sn += 1
    add_bullet_slide(prs, "Section 6 - Data Types", [
        "6.4   Integer types: 2-state (bit, int) vs 4-state (logic, integer)",
        "6.5   wire vs reg vs logic: where each is legal / illegal",
        "6.18  User-defined types: typedef, const, parameter, localparam",
        "6.19  Enumerations: named constants with methods",
        "6.24  Type casting: static cast type'() and dynamic $cast()",
        "6.20  Type parameters: parameter type for generic modules",
        "6.25  Parameterized data types: generic containers via classes",
    ], sn)

    sn += 1
    add_code_slide(prs, "6.4 Integer Data Types", """
// 2-state types: 0, 1 only (default = 0)
bit       b;               // 1-bit
byte      by = 127;        // signed 8-bit
int       i;               // signed 32-bit
longint   li;              // signed 64-bit

$display("bit=%b, byte=%0d, int=%0d", b, by, i);
by = by + 1;
$display("byte 127+1 = %0d (overflow wraps)", by);

// 4-state types: 0, 1, x, z (default = x)
logic      l;              // 1-bit
integer    ig;             // signed 32-bit
logic [7:0] l8;

$display("logic=%b, integer=%b (x=unknown)", l, ig);

l8 = 8'bxxxx_zzzz;
$display("logic[7:0] = %b (x and z)", l8);
""",
    "bit=0, byte=127, int=0\n"
    "byte 127+1 = -128 (overflow wraps)\n"
    "logic=x, integer=x (default)\n"
    "logic[7:0] = xxxxzzzz", sn)

    sn += 1
    add_code_slide(prs, "6.5 wire vs reg vs logic", """
//            assign  always_ff  always_comb  module port
// wire:       YES      NO          NO        input/output
// reg:        NO       YES         YES       output only
// logic:      YES      YES         YES       input/output

// wire: ONLY with continuous assign or port connections
wire [7:0] w;
assign w = 8'hAA;         // OK
// always @(*) w = 8'hBB; // ILLEGAL: wire in procedural

// reg: ONLY in procedural blocks (legacy Verilog)
reg [7:0] r;
// assign r = 8'hBB;      // ILLEGAL: reg in assign
always_comb r = 8'hBB;    // OK

// logic: works EVERYWHERE (SystemVerilog replacement)
logic [7:0] l;
assign l = 8'hCC;         // OK: continuous assign
// always_ff @(posedge clk) l <= 8'hDD; // also OK

$display("wire=%h, reg=%h, logic=%h", w, r, l);
$display("Use logic everywhere (replaces reg)");
$display("Use wire only when multiple drivers needed");
""",
    "wire=aa, reg=bb, logic=cc\n"
    "logic replaces reg in SystemVerilog\n"
    "wire needed only for multi-driver buses", sn)

    sn += 1
    add_code_slide(prs, "6.19 Enumerations", """
typedef enum {IDLE, READ, WRITE, ERROR} state_t;

typedef enum logic [2:0] {
    CMD_NOP  = 3'b000,
    CMD_READ = 3'b010,
    CMD_WRITE= 3'b011,
    CMD_RST  = 3'b111
} cmd_t;

state_t st = IDLE;
$display("state = %s (value=%0d)", st.name(), st);

// Enum methods
st = st.first(); $display("first = %s", st.name());
st = st.last();  $display("last  = %s", st.name());
st = IDLE;
st = st.next();  $display("IDLE.next = %s", st.name());

// Iterate all values
st = st.first();
do begin
    $display("  %s = %0d", st.name(), st);
    st = st.next();
end while (st != st.first());
""",
    "state = IDLE (value=0)\n"
    "first=IDLE, last=ERROR\n"
    "IDLE.next = READ\n"
    "IDLE=0, READ=1, WRITE=2, ERROR=3", sn)

    sn += 1
    add_code_slide(prs, "6.18 typedef & Constants", """
// typedef: create named types
typedef logic [7:0]  byte_t;
typedef logic [31:0] word_t;

byte_t b = 8'hAB;
word_t w = 32'hDEAD_BEEF;
$display("byte_t = %h, word_t = %h", b, w);

// parameter: overridable at instantiation
parameter int DATA_WIDTH = 8;
// localparam: computed, cannot be overridden
localparam int ADDR_BITS = $clog2(256);
$display("DATA_WIDTH=%0d, ADDR_BITS=%0d", DATA_WIDTH, ADDR_BITS);

// const: runtime constant
const int VERSION = 42;
$display("VERSION = %0d (const)", VERSION);

// parameter vs localparam:
// parameter  -> can be overridden (#(.WIDTH(16)))
// localparam -> derived, cannot be overridden
""",
    "byte_t = ab, word_t = deadbeef\n"
    "DATA_WIDTH=8, ADDR_BITS=8\n"
    "VERSION = 42", sn)

    sn += 1
    add_code_slide(prs, "6.24 Casting & Scope Resolution", """
// Static cast: type'(expression)
int i = 42;
real r = real'(i);
$display("int->real: %f", r);

r = 3.99;
i = int'(r);
$display("real->int: %0d (truncated)", i);

logic [7:0] v = 8'hFF;
$display("unsigned: %0d", v);
$display("signed:   %0d", $signed(v));

// $cast: dynamic (runtime) check
typedef enum {RED, GREEN, BLUE} color_t;
color_t c;
if ($cast(c, 1))
    $display("$cast(1) = %s", c.name());

// Scope resolution ::
// math_pkg::max(a, b)
// class_name::static_method()
""",
    "int->real: 42.000000\n"
    "real->int: 3 (truncated)\n"
    "unsigned: 255, signed: -1\n"
    "$cast(1) = GREEN", sn)

    sn += 1
    add_code_slide(prs, "6.20 Parameterized Types - Type Parameters", """
// Module with a type parameter (LRM 6.20.3)
module generic_register #(parameter type T = logic [7:0])(
    input  logic clk, rst, en,
    input  T d, output T q
);
    always_ff @(posedge clk or posedge rst)
        if (rst) q <= T'(0);
        else if (en) q <= d;
endmodule

// Testbench: instantiate with different types
logic [7:0] d8  = 8'hA5, q8;    // 8-bit (default)
int         d32 = 123456, q32;   // 32-bit integer

generic_register            u1 (
    .clk(clk), .rst(rst), .en(1'b1), .d(d8),  .q(q8));
generic_register #(.T(int)) u2 (
    .clk(clk), .rst(rst), .en(1'b1), .d(d32), .q(q32));

$display("8-bit reg:  q8  = %h", q8);   // d8=A5
$display("32-bit reg: q32 = %0d", q32);  // d32=123456
""",
    "8-bit reg:  q8  = a5\n"
    "32-bit reg: q32 = 123456\n"
    "parameter type T: one module works for any type", sn)

    sn += 1
    add_code_slide(prs, "6.25 Parameterized Data Types", """
// Virtual class as a type container (LRM 6.25)
virtual class C #(parameter type T = logic,
                  parameter int SIZE = 1);
    typedef T [SIZE-1:0] t_vector;
    typedef T            t_element;
endclass

// Specialize with different types
typedef C#(logic, 8)::t_vector byte_vec_t;
typedef C#(bit,   4)::t_vector nibble_t;

byte_vec_t bv = 8'hAB;
nibble_t   nb = 4'b1010;
$display("byte_vec_t = %h", bv);
$display("nibble_t   = %b", nb);

// Real-world use: generic FIFO types, bus widths, etc.
// One class definition, many specializations
$display("Parameterized types avoid code duplication");
""",
    "byte_vec_t = ab\nnibble_t   = 1010\n"
    "Parameterized types avoid code duplication", sn)

    # ================================================================
    #  LAB: SECTION 6
    # ================================================================
    sn += 1
    add_bullet_slide(prs, "Lab Exercise - Data Types", [
        "Declare an enum: typedef enum {RED, GREEN, BLUE} color_t;",
        "Use .name() and .next() to print all values with $display.",
    ], sn)

    # ================================================================
    #  SECTION 3: DESIGN BUILDING BLOCKS
    # ================================================================
    sn += 1
    add_bullet_slide(prs, "Section 3 - Design Building Blocks", [
        "3.2   Modules: basic design unit, ports, parameterization",
        "3.10  Packages: reusable definitions, import, scope resolution",
        "3.12  Interfaces: bundled signals with modports",
        "3.13  Subroutines: functions (no delay) and tasks (with delay)",
        "3.14  Timeunits: timeunit and timeprecision declarations",
    ], sn)

    sn += 1
    add_code_slide(prs, "3.2 Modules - Declaration & Instantiation", """
module adder (
    input  logic [7:0] a, b,
    output logic [8:0] sum
);
    assign sum = a + b;
endmodule

module counter #(parameter WIDTH = 8)(
    input  logic             clk, rst,
    output logic [WIDTH-1:0] count
);
    always_ff @(posedge clk or posedge rst)
        if (rst) count <= '0;
        else     count <= count + 1;
endmodule

// Named port instantiation (explicit mapping):
adder u_add (.a(op1), .b(op2), .sum(result));

// Parameterized instantiation:
counter #(.WIDTH(4)) u_cnt4 (
    .clk(clk), .rst(rst), .count(cnt4)
);
""",
    "Modules are the basic building blocks\n"
    "Ports: input, output, inout\n"
    "Parameters override with #(.NAME(val))", sn)

    sn += 1
    add_code_slide(prs, "3.12 Interfaces - Signals & Modports", """
interface bus_if;
    logic       valid;
    logic       ready;
    logic [7:0] data;

    modport master (output valid, data, input  ready);
    modport slave  (input  valid, data, output ready);
endinterface

module producer (bus_if.master bus);
    // master can DRIVE valid and data
    // master can READ ready
endmodule

module consumer (bus_if.slave bus);
    // slave can READ valid and data
    // slave can DRIVE ready
endmodule

// Without interface: 3 separate ports per module
// With interface:    1 bundled port per module
// Modports enforce direction at compile time
""",
    "Interface: bundles related signals into one port\n"
    "Modport: controls direction per module\n"
    "master: drives valid/data, reads ready", sn)

    sn += 1
    add_code_slide(prs, "3.13 Functions & Tasks", """
// Functions: no delays, return a value
function automatic int factorial(int n);
    if (n <= 1) return 1;
    else        return n * factorial(n - 1);
endfunction

// Tasks: CAN have delays, no return value
task automatic delayed_print(string msg, int dly);
    #dly;
    $display("[%0t] %s", $time, msg);
endtask

// Pass by reference
function automatic void swap(ref int a, ref int b);
    int t; t = a; a = b; b = t;
endfunction

$display("factorial(5) = %0d", factorial(5));
$display("factorial(0) = %0d", factorial(0));
""",
    "factorial(5) = 120\nfactorial(0) = 1\n"
    "Functions: combinational, Tasks: can have #delay", sn)

    sn += 1
    add_code_slide(prs, "3.10 Packages - Reusable Definitions", """
package math_pkg;
    parameter int PI_INT = 3;
    typedef struct { int x; int y; } point_t;
    function automatic int max(int a, int b);
        return (a > b) ? a : b;
    endfunction
endpackage

module top;
    import math_pkg::point_t;  // explicit import
    import math_pkg::max;

    point_t p;
    p.x = 10; p.y = 20;
    $display("max(%0d,%0d) = %0d", p.x, p.y, max(p.x, p.y));

    // Scope resolution operator ::
    $display("PI = %0d", math_pkg::PI_INT);
endmodule
""",
    "max(10,20) = 20\nPI = 3", sn)

    sn += 1
    add_code_slide(prs, "3.14 Timeunits", """
module timing_demo;
    timeunit 1ns;
    timeprecision 1ps;

    initial begin
        $timeformat(-9, 3, " ns", 12);

        $display("Time at start: %t", $time);
        #1;
        $display("After #1:      %t", $time);
        #0.5;
        $display("After #0.5:    %t", $time);
        #10.123;
        $display("After #10.123: %t", $time);

        $display("timeunit = 1ns");
        $display("timeprecision = 1ps");
    end
endmodule
""",
    "Time at start:    0.000 ns\n"
    "After #1:         1.000 ns\n"
    "After #0.5:       1.500 ns\n"
    "After #10.123:   11.623 ns", sn)

    # ================================================================
    #  SECTION 4: SCHEDULING SEMANTICS
    # ================================================================
    sn += 1
    add_bullet_slide(prs, "Section 4 - Scheduling Semantics", [
        "SystemVerilog simulation is event-driven",
        "Blocking (=): executes sequentially, sees updated values immediately",
        "Non-blocking (<=): samples BEFORE updating (all RHS read first)",
        "Rule: use = in combinational logic (always_comb)",
        "Rule: use <= in sequential logic (always_ff)",
        "Never mix = and <= for the same signal",
        "Race conditions occur when multiple blocks drive the same signal with =",
    ], sn)

    sn += 1
    add_code_slide(prs, "4.5 Blocking (=) vs Non-Blocking (<=)", """
// Blocking: sequential execution
a = 1;
b = a;   // sees a=1
c = b;   // sees b=1
$display("Blocking:     a=%b b=%b c=%b", a, b, c);

// Non-blocking: all RHS sampled FIRST, then updated
x <= 0;
y <= x;  // samples OLD x (=1), schedules y=1
z <= y;  // samples OLD y (=0), schedules z=0
#1;
$display("Non-Blocking: x=%b y=%b z=%b", x, y, z);

// Pipeline (correct with <=)
always @(posedge clk) pipe_a <= 8'hAA;
always @(posedge clk) pipe_b <= pipe_a; // gets OLD pipe_a
""",
    "Blocking:     a=1 b=1 c=1\n"
    "Non-Blocking: x=0 y=1 z=0\n"
    "Use = in comb, <= in sequential logic", sn)

    # ================================================================
    #  LAB: SECTION 4
    # ================================================================
    sn += 1
    add_bullet_slide(prs, "Lab Exercise - Scheduling Semantics", [
        "Write two always blocks that both assign to the same signal 'y'",
        "using blocking (=) at the same simulation time.",
        "Run the simulation multiple times and observe whether the value",
        "of 'y' changes between runs. This is a race condition.",
    ], sn)

    # ================================================================
    #  SECTION 7: AGGREGATE DATA TYPES
    # ================================================================
    sn += 1
    add_bullet_slide(prs, "Section 7 - Aggregate Data Types", [
        "7.2  Structures (packed and unpacked)",
        "7.3  Unions (packed and tagged)",
        "7.4  Packed and Unpacked Arrays",
        "7.5  Dynamic Arrays",
        "7.6  Array Assignments",
        "7.7  Arrays as Arguments to Subroutines",
        "7.8  Associative Arrays",
        "7.9  Associative Array Methods",
        "7.10 Queues",
        "7.12 Array Manipulation Methods",
    ], sn)

    # ===== 7.2 STRUCTURES =====
    sn += 1
    add_code_slide(prs, "7.2 Structures - Basic Declaration", """
struct { bit [7:0] opcode; bit [23:0] addr; } IR;
IR.opcode = 8'hAB;
IR.addr   = 24'h123456;
$display("IR.opcode = %h", IR.opcode);
$display("IR.addr   = %h", IR.addr);

typedef struct {
    bit [7:0]  opcode;
    bit [23:0] addr;
} instruction;

instruction IR2;
IR2.opcode = 8'hFF;
$display("IR2.opcode = %h", IR2.opcode);
""",
    "IR.opcode = ab\nIR.addr   = 123456\nIR2.opcode = ff", sn)

    sn += 1
    add_code_slide(prs, "7.2.1 Packed Structures", """
struct packed signed {
    int       a;      // bits [63:32]
    shortint  b;      // bits [31:16]
    byte      c;      // bits [15:8]
    bit [7:0] d;      // bits [7:0]
} pack1;

pack1.a = 1; pack1.b = 2; pack1.c = 3; pack1.d = 4;
$display("pack1 as vector = %0h", pack1);
$display("pack1[15:8] = %0h", pack1[15:8]); // same as c
""",
    "pack1 as vector = 100020304\npack1[15:8] = 3 (same as pack1.c)", sn)

    sn += 1
    add_code_slide(prs, "7.2.1 Packed Struct - ATM Header", """
typedef struct packed {
    bit [3:0]  GFC;
    bit [7:0]  VPI;
    bit [11:0] VCI;
    bit        CLP;
    bit [3:0]  PT;
    bit [7:0]  HEC;
} s_atmheader;

s_atmheader hdr;
hdr = '{GFC:4'hA, VPI:8'hBB, VCI:12'hCCC,
        CLP:1'b1, PT:4'h5, HEC:8'hFF};
$display("GFC=%h VPI=%h VCI=%h CLP=%b PT=%h HEC=%h",
         hdr.GFC, hdr.VPI, hdr.VCI, hdr.CLP, hdr.PT, hdr.HEC);
$display("Bit width = %0d bits", $bits(hdr));
""",
    "GFC=a VPI=bb VCI=ccc CLP=1 PT=5 HEC=ff\n"
    "Bit width = 38 bits", sn)

    sn += 1
    add_code_slide(prs, "7.2.2 Assigning to Structures", """
typedef struct {
    int  addr;
    int  crc;
    byte data [4];
} packet1;

packet1 p1;

// Assignment pattern (positional)
p1 = '{10, 32'hDEAD, '{8'hAA, 8'hBB, 8'hCC, 8'hDD}};
$display("p1.addr = %0d, p1.crc = %h", p1.addr, p1.crc);

// Named assignment pattern
p1 = '{addr: 42, crc: 32'hBEEF, data: '{1, 2, 3, 4}};
$display("p1.addr = %0d, p1.crc = %h", p1.addr, p1.crc);
""",
    "p1.addr = 10, p1.crc = dead\np1.addr = 42, p1.crc = beef", sn)

    # ===== 7.3 UNIONS =====
    sn += 1
    add_code_slide(prs, "7.3 Unions - Basic & Packed", """
typedef union { int i; shortreal f; } num;
num n;
n.i = 42;       $display("n.i = %0d", n.i);
n.f = 3.14;     $display("n.f = %f", n.f);

typedef union packed {
    logic [31:0]     word;
    logic [3:0][7:0] bytes;
    logic [1:0][15:0] halves;
} word_u;

word_u wu;
wu.word = 32'hDEAD_BEEF;
$display("wu.bytes[3] = %h", wu.bytes[3]);
$display("wu.halves[0] = %h", wu.halves[0]);
""",
    "n.i = 42\nn.f = 3.14\n"
    "wu.bytes[3] = de (MSB)\nwu.halves[0] = beef (lower)", sn)

    # ===== 7.4 ARRAYS =====
    sn += 1
    add_code_slide(prs, "7.4 Packed vs Unpacked Arrays", """
bit [7:0] c1;            // packed array (vector)
real u [7:0];            // unpacked array

bit [3:0][7:0] data32;   // 4 packed bytes = 32 bits
data32 = 32'hCAFE_BABE;
$display("data32 = %h", data32);
$display("data32[3] = %h, data32[0] = %h",
         data32[3], data32[0]);
$display("data32 + 1 = %h", data32 + 1);

int Array[0:7][0:31];    // unpacked using ranges
int Array2[8][32];       // same thing using sizes
""",
    "data32 = cafebabe\n"
    "data32[3] = ca, data32[0] = be\ndata32 + 1 = cafebabf", sn)

    sn += 1
    add_code_slide(prs, "7.4.4 Memories & 7.4.5 Multidimensional", """
// Memory (1D array of reg/logic/bit)
logic [7:0] mema [0:255];
mema[5] = 8'hAA;
$display("mema[5] = %h", mema[5]);

// Multidimensional: 10 elements of 4 packed bytes
bit [3:0][7:0] joe [1:10];
joe[1] = 32'h01020304;
joe[3] = joe[1] + joe[2];
$display("joe[3] = %h", joe[3]);
joe[4][3:2] = joe[1][1:0];
$display("joe[4][3]=%h, joe[4][2]=%h",
         joe[4][3], joe[4][2]);
""",
    "mema[5] = aa\n"
    "joe[3] = 06080a0c\n"
    "joe[4][3]=03, joe[4][2]=04", sn)

    sn += 1
    add_code_slide(prs, "7.4.6 Indexing, Slicing, and Part-Select", """
logic [63:0] data = 64'h0123_4567_89AB_CDEF;
$display("data[23:16] = %h", data[23:16]);

bit [3:0][7:0] j = 32'hAABBCCDD;
$display("j[2] = %h", j[2]);

bit signed [31:0] busA [7:0];
busB = busA[7:6];              // 2-element slice

// Variable-width part-select
logic [31:0] bitvec = 32'hDCBA_9876;
$display("bitvec[8+:8] = %h", bitvec[8 +: 8]);
$display("bitvec[23-:8] = %h", bitvec[23 -: 8]);
""",
    "data[23:16] = ab\nj[2] = bb\n"
    "bitvec[8+:8] = 98\nbitvec[23-:8] = ba", sn)

    sn += 1
    add_code_slide(prs, "7.11 Array Querying Functions", """
logic [7:0] arr [3:0][2:0];

$display("$left=%0d, $right=%0d", $left(arr,1), $right(arr,1));
$display("$low=%0d,  $high=%0d",  $low(arr,1),  $high(arr,1));
$display("$size(1)=%0d, $size(2)=%0d", $size(arr,1), $size(arr,2));
$display("$dimensions=%0d", $dimensions(arr));
$display("$unpacked_dimensions=%0d", $unpacked_dimensions(arr));
$display("$increment=%0d", $increment(arr, 1));
""",
    "$left=3, $right=0\n$low=0,  $high=3\n"
    "$size(1)=4, $size(2)=3\n"
    "$dimensions=3, $unpacked_dimensions=2\n"
    "$increment=-1 (descending)", sn)

    # ===== 7.5 DYNAMIC ARRAYS =====
    sn += 1
    add_code_slide(prs, "7.5 Dynamic Arrays", """
int arr[];
arr = new[5];
$display("arr.size() = %0d", arr.size());

int isrc[3] = '{5, 6, 7};
int idest[];
idest = new[3](isrc);
foreach (idest[i]) $display("idest[%0d] = %0d", i, idest[i]);

// Resize: truncate or pad
int src[3] = '{2, 3, 4};
int dest1[];
dest1 = new[2](src);       // '{2, 3}       (truncated)
int dest2[];
dest2 = new[5](src);       // '{2, 3, 4, 0, 0} (padded)
""",
    "arr.size() = 5\nidest = '{5, 6, 7}\n"
    "dest1 = '{2, 3}\ndest2 = '{2, 3, 4, 0, 0}", sn)

    sn += 1
    add_code_slide(prs, "7.5.2-7.5.3 size() and delete()", """
int arr[];
arr = new[10];
$display("size = %0d", arr.size());  // 10

// Quadruple
arr = new[arr.size() * 4](arr);
$display("size = %0d", arr.size());  // 40

// Delete all elements
int ab[] = new[5];
ab.delete;
$display("size = %0d", ab.size());   // 0
""",
    "size = 10\nsize = 40\nsize = 0", sn)

    # ===== 7.6 ARRAY ASSIGNMENTS =====
    sn += 1
    add_code_slide(prs, "7.6 Array Assignments", """
int A[10:1], B[0:9];
A = B;             // OK: same type and size (10 elements)
$display("A[10]=%0d, A[9]=%0d", A[10], A[9]);

int C[];
C = A;             // Dynamic auto-resizes
$display("C.size() = %0d", C.size());

// Unpacked array concatenation
string d[1:5] = '{"a","b","c","d","e"};
string p[];
p = {d[1:3], "hello", d[4:5]};
foreach (p[i]) $display("p[%0d] = %s", i, p[i]);
""",
    "A[10]=B[0], A[9]=B[1]\n"
    "C.size() = 10\np = '{a, b, c, hello, d, e}", sn)

    # ===== 7.7 ARRAYS AS ARGUMENTS =====
    sn += 1
    add_code_slide(prs, "7.7 Arrays as Arguments to Subroutines", """
task print_matrix(int a[3:1][3:1]);
    foreach (a[i,j]) $display("a[%0d][%0d]=%0d", i, j, a[i][j]);
endtask

int b[3:1][3:1];   // OK: same type, dimension, size
int b[1:3][0:2];   // OK: same size (different ranges)
// int b[3:1];      // ERROR: wrong number of dimensions
// int b[3:1][4:1]; // ERROR: wrong size (3 vs 4)

// Dynamic can be passed to fixed-size formal
task t(string arr[4:1]);
    foreach (arr[i]) $display("arr[%0d]=%s", i, arr[i]);
endtask

// Pass by reference modifies caller's array
task double_elements(ref int arr[]);
    foreach (arr[i]) arr[i] = arr[i] * 2;
endtask
""",
    "Compatible: same type + dimension + size\n"
    "Dynamic->Fixed: runtime size check\n"
    "ref: modifications visible to caller", sn)

    # ===== 7.8 ASSOCIATIVE ARRAYS =====
    sn += 1
    add_code_slide(prs, "7.8 Associative Arrays", """
integer i_array[*];          // wildcard index
bit [20:0] arr_b[string];   // string index
int scores[int];             // integral index

// String-indexed
arr_b["hello"] = 21'h1AAAAA;
$display("arr_b[hello] = %h", arr_b["hello"]);

// Int-indexed (signed ordering)
scores[-10] = 50; scores[0] = 75; scores[10] = 100;
foreach (scores[k])
    $display("scores[%0d] = %0d", k, scores[k]);

// Auto-allocation on assignment (7.8.7)
int a[int] = '{default:1};
a[1]++;
$display("a[1] = %0d", a[1]);
""",
    "arr_b[hello] = 1aaaaa\n"
    "scores[-10]=50, scores[0]=75, scores[10]=100\n"
    "a[1] = 2 (default 1 + increment)", sn)

    # ===== 7.9 ASSOCIATIVE ARRAY METHODS =====
    sn += 1
    add_code_slide(prs, "7.9 Associative Array Methods", """
int map[string];
map["hello"]=1; map["sad"]=2; map["world"]=3;

$display("num = %0d, size = %0d", map.num(), map.size());
map.delete("sad");
$display("After delete(sad): size = %0d", map.size());
map.delete;
$display("After delete all: size = %0d", map.size());

map["test"] = 99;
$display("exists(test) = %0d", map.exists("test"));
$display("exists(nope) = %0d", map.exists("nope"));
""",
    "num = 3, size = 3\n"
    "After delete(sad): size = 2\n"
    "After delete all: size = 0\n"
    "exists(test) = 1, exists(nope) = 0", sn)

    sn += 1
    add_code_slide(prs, "7.9 Traversal: first/last/next/prev", """
int map[string];
map["banana"]=2; map["apple"]=1;
map["cherry"]=3; map["date"]=4;
string s;

// Forward traversal (lexicographic)
if (map.first(s))
    do $display("map[%s]=%0d", s, map[s]);
    while (map.next(s));

// Reverse traversal
if (map.last(s))
    do $display("map[%s]=%0d", s, map[s]);
    while (map.prev(s));
""",
    "Forward:  apple=1, banana=2, cherry=3, date=4\n"
    "Reverse:  date=4, cherry=3, banana=2, apple=1", sn)

    sn += 1
    add_code_slide(prs, "7.9.11 Associative Array Literals", """
// Literal with default value
string words[int] = '{default: "hello"};
$display("words[999] = %s", words[999]);

// Literal with entries and default
integer tab[string] = '{
    "Peter": 20,
    "Paul":  22,
    "Mary":  23,
    default: -1
};

$display("tab[Peter] = %0d", tab["Peter"]);
$display("tab[unknown] = %0d", tab["unknown"]);
""",
    "words[999] = hello\n"
    "tab[Peter] = 20\n"
    "tab[unknown] = -1 (default)", sn)

    # ===== 7.10 QUEUES =====
    sn += 1
    add_code_slide(prs, "7.10 Queues - Declaration & Operators", """
byte    q1[$];                // unbounded queue
string  names[$] = {"Bob"};   // initialized
integer Q[$] = {3, 2, 7};    // initialized
bit     q2[$:255];            // bounded (max 256)

$display("q1.size() = %0d", q1.size());
$display("names[0] = %s", names[0]);
foreach (Q[i]) $display("Q[%0d] = %0d", i, Q[i]);

// Slicing
$display("Q[0:$] = entire queue");
$display("Q[1:2] slice:");
foreach (Q[i]) if (i>=1 && i<=2) $display("  [%0d]=%0d", i, Q[i]);
""",
    "q1.size() = 0, names[0] = Bob\n"
    "Q = '{3, 2, 7}", sn)

    sn += 1
    add_code_slide(prs, "7.10.2 Queue Methods", """
int Q[$] = {10, 20, 30};
int e;

$display("size = %0d", Q.size());
Q.push_front(5);     $display("push_front(5): Q=%p", Q);
Q.push_back(40);     $display("push_back(40): Q=%p", Q);
e = Q.pop_front();   $display("pop_front = %0d", e);
e = Q.pop_back();    $display("pop_back = %0d", e);
Q.insert(1, 15);     $display("insert(1,15): Q=%p", Q);
Q.delete(2);         $display("delete(2): Q=%p", Q);
""",
    "size = 3\npush_front(5): '{5,10,20,30}\n"
    "pop_front=5, pop_back=40\n"
    "insert(1,15): '{10,15,20,30}\n"
    "delete(2): '{10,15,30}", sn)

    sn += 1
    add_code_slide(prs, "7.10.4 Queue via Assignment", """
int q[$] = {2, 4, 8};
int e = 1;

q = {q, 6};        $display("push_back(6):  %p", q);
q = {e, q};        $display("push_front(1): %p", q);
q = q[1:$];        $display("pop_front:     %p", q);
q = q[0:$-1];      $display("pop_back:      %p", q);
q = {};            $display("clear: size = %0d", q.size());

// Equivalent to insert(pos, e)
// q = {q[0:pos-1], e, q[pos:$]};

// Advanced slices
// q = q[2:$];      // remove first two items
// q = q[1:$-1];    // remove first and last
""",
    "push_back(6):  '{2,4,8,6}\n"
    "push_front(1): '{1,2,4,8,6}\n"
    "pop_front:     '{2,4,8,6}\n"
    "pop_back:      '{2,4,8}\n"
    "clear: size = 0", sn)

    # ===== 7.12 ARRAY MANIPULATION =====
    sn += 1
    add_code_slide(prs, "7.12.1 Array Locator Methods", """
int arr[] = '{10, 20, 30, 20, 40};
int qi[$];

qi = arr.find(x) with (x > 15);
$display("find(>15) = %p", qi);
qi = arr.find_index with (item == 20);
$display("find_index(==20) = %p", qi);

qi = arr.min;    $display("min = %p", qi);
qi = arr.max;    $display("max = %p", qi);
qi = arr.unique; $display("unique = %p", qi);
""",
    "find(>15) = '{20,30,20,40}\n"
    "find_index(==20) = '{1,3}\n"
    "min='{10}, max='{40}, unique='{10,20,30,40}", sn)

    sn += 1
    add_code_slide(prs, "7.12.2 Array Ordering Methods", """
string s[] = '{"hello", "sad", "world"};
s.reverse;    $display("reverse: %p", s);

int q[$] = '{4, 5, 3, 1};
q.sort;       $display("sort:  %p", q);
q.rsort;      $display("rsort: %p", q);
q.shuffle;    $display("shuffle: %p", q);

// Sort structs by field
struct { byte red, green, blue; } c[512];
c.sort with (item.red);
c.sort(x) with ({x.blue, x.green});
""",
    "reverse: '{world, sad, hello}\n"
    "sort:  '{1, 3, 4, 5}\n"
    "rsort: '{5, 4, 3, 1}\n"
    "shuffle: (random order)", sn)

    sn += 1
    add_code_slide(prs, "7.12.3 Array Reduction Methods", """
byte b[] = '{1, 2, 3, 4};

$display("sum = %0d", b.sum);
$display("product = %0d", b.product);
$display("and = %h", b.and);
$display("or  = %h", b.or);
$display("xor = %h", b.xor);
$display("xor(item+4) = %0d", b.xor with (item + 4));

// 2D sum
logic [7:0] m[2][2] = '{'{ 5,10}, '{15,20}};
$display("2D sum = %0d",
         m.sum with (item.sum with (item)));
""",
    "sum=10, product=24\nand=00, or=07, xor=04\n"
    "xor(item+4) = 12\n2D sum = 50", sn)

    sn += 1
    add_code_slide(prs, "7.12.4 Iterator Index Querying", """
int arr[] = '{0, 10, 2, 30, 4};
int qi[$];

// Find items equal to their index position
qi = arr.find with (item == item.index);
$display("Items equal to index: %p", qi);
// Index 0: value=0  (0==0 YES)
// Index 1: value=10 (10==1 NO)
// Index 2: value=2  (2==2 YES)
// Index 3: value=30 (30==3 NO)
// Index 4: value=4  (4==4 YES)

qi = arr.find_index with (item == item.index);
$display("Matching indices: %p", qi);
""",
    "Items equal to index: '{0, 2, 4}\n"
    "Matching indices: '{0, 2, 4}", sn)

    # ================================================================
    #  LAB: SECTION 7
    # ================================================================
    sn += 1
    add_bullet_slide(prs, "Lab Exercise - Aggregate Data Types", [
        "Create a queue of integers, push values 1 to 5 using push_back(),",
        "pop the first two using pop_front(), and display the remaining items.",
    ], sn)

    # ================================================================
    #  SUMMARY SLIDE
    # ================================================================
    sn += 1
    add_bullet_slide(prs, "Summary", [
        "Section 3: Modules, interfaces, packages, functions/tasks, programs",
        "Section 4: Blocking (=) for comb, non-blocking (<=) for sequential",
        "Section 5: Number formats, operators, string literals",
        "Section 6: 2-state vs 4-state, enums, typedef, parameterized types",
        "Section 7: Structs, unions, arrays (fixed, dynamic, associative, queues)",
        "Array methods: find, sort, sum, unique with 'with' clause",
        "Parameterized types: parameter type T for generic reusable code",
        "QuestaSim: vlib, vlog -sv, vsim, run -all, .do scripts",
    ], sn)

    os.makedirs(OUTPUT_DIR, exist_ok=True)
    prs.save(OUTPUT_FILE)
    print(f"Presentation saved to: {OUTPUT_FILE}")
    print(f"Total slides: {sn}")


if __name__ == "__main__":
    main()
